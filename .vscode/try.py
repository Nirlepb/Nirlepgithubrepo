from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
layout = prs.slide_layouts[6]

# Intro slide
slide = prs.slides.add_slide(layout)
title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
title_tf = title_box.text_frame
title = title_tf.add_paragraph()
title.text = "🌸 The Art and Science of Meditation 🌸"
title.font.size = Pt(54)
title.font.bold = True
title.font.color.rgb = RGBColor(240, 210, 150)
title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
subtitle_tf = subtitle_box.text_frame
subtitle = subtitle_tf.add_paragraph()
subtitle.text = "A Journey of Mindfulness, Peace, and Science"
subtitle.font.size = Pt(28)
subtitle.font.color.rgb = RGBColor(220, 220, 220)
subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slides
sections = [
    ("Ancient Wisdom Meets Modern Understanding", "Meditation has been practiced for thousands of years across cultures and spiritual traditions. It offers universal principles for calming the mind and connecting with our deeper selves."),
    ("How Meditation Transforms Your Brain", "Neuroscience shows that meditation strengthens the prefrontal cortex and reduces amygdala activity, enhancing emotional regulation and calm."),
    ("The Profound Benefits of Meditation", "Meditation improves mental clarity, emotional balance, physical health, and inner peace."),
    ("Exploring Different Meditation Practices", "Zen, Mindfulness, Loving-Kindness, and Mantra meditation each offer unique ways to connect with the present moment."),
    ("Creating Your Meditation Space", "A calm corner with soft lighting, plants, and gentle sounds enhances your meditation experience."),
    ("Building Your Daily Meditation Habit", "Start small, stay consistent, and let your progress naturally unfold."),
    ("Overcoming Common Meditation Challenges", "Distractions are part of the process. Gently return focus to your breath and maintain patience."),
    ("Your Journey Toward Inner Peace Begins Now", "Each moment of mindfulness is a step closer to calm and clarity. Begin your practice today.")
]

for title, content in sections:
    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.add_paragraph()
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(235, 210, 160)

    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4))
    content_tf = content_box.text_frame
    p2 = content_tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(230, 230, 230)

prs.save("The_Art_and_Science_of_Meditation_Enhanced_Template.pptx")
print("✅ PowerPoint saved successfully!")
